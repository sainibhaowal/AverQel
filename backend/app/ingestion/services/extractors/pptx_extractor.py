from __future__ import annotations

from collections.abc import Callable
from io import BytesIO

from app.core.errors import ApiError
from app.ingestion.services.extractors.base import (
    BaseExtractor,
    ExtractionRequest,
    ExtractionResult,
)
from app.ingestion.services.parser_service import sanitize_document_text


class PptxExtractor(BaseExtractor):
    extraction_method = "pptx_native"
    supported_extensions = frozenset({".pptx"})
    supported_mime_types = frozenset(
        {"application/vnd.openxmlformats-officedocument.presentationml.presentation"}
    )

    def __init__(self, *, max_text_chars: int) -> None:
        self.max_text_chars = max_text_chars

    def extract(self, request: ExtractionRequest) -> ExtractionResult:
        constructor = self._load_presentation_constructor()
        try:
            presentation = constructor(BytesIO(request.payload))
        except Exception as exc:  # noqa: BLE001
            raise ApiError(
                code="PPTX_PARSE_FAILED",
                message="Failed to parse PPTX presentation.",
                status_code=422,
            ) from exc

        lines: list[str] = []
        slide_count = 0
        for slide_index, slide in enumerate(
            getattr(presentation, "slides", []), start=1
        ):
            slide_count += 1
            slide_lines: list[str] = []
            for shape in getattr(slide, "shapes", []):
                text = sanitize_document_text(getattr(shape, "text", "")).strip()
                if text:
                    slide_lines.append(text)

            notes = ""
            notes_slide = getattr(slide, "notes_slide", None)
            if notes_slide is not None:
                notes_text_frame = getattr(notes_slide, "notes_text_frame", None)
                notes = sanitize_document_text(
                    getattr(notes_text_frame, "text", "")
                ).strip()
                if notes:
                    slide_lines.append(f"[notes] {notes}")

            if slide_lines:
                lines.append(f"## Slide {slide_index}\n" + "\n".join(slide_lines))

        text = "\n\n".join(lines).strip()
        if len(text) > self.max_text_chars:
            raise ApiError(
                code="DOCUMENT_TEXT_LIMIT_EXCEEDED",
                message="Parsed document exceeds text processing limit.",
                status_code=422,
                details={"max_chars": self.max_text_chars},
            )

        warnings: list[str] = []
        if not text:
            warnings.append("pptx_no_text_extracted")

        return ExtractionResult(
            text=text,
            page_count=slide_count,
            extraction_method=self.extraction_method,
            coverage_score=1.0 if text else 0.0,
            warnings=warnings,
        )

    def _load_presentation_constructor(self) -> Callable[[BytesIO], object]:
        try:
            from pptx import Presentation
        except Exception as exc:  # noqa: BLE001
            raise ApiError(
                code="PPTX_EXTRACTOR_DEPENDENCY_MISSING",
                message="python-pptx is required for PPTX extraction.",
                status_code=503,
            ) from exc
        return Presentation
