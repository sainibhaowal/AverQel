#!/usr/bin/env python3
import shutil
from pathlib import Path

# Need docx, pptx, openpyxl, reportlab, Pillow
# They should all be in the environment!
try:
    from docx import Document
    from openpyxl import Workbook  # type: ignore[import-untyped]
    from PIL import Image, ImageDraw, ImageFilter
    from pptx import Presentation
    from reportlab.pdfgen import canvas  # type: ignore[import-untyped]
except ImportError as e:
    print(f"Missing dependency: {e}")
    exit(1)

DATASET_DIR = Path(__file__).parent.parent.parent / "Docs" / "Datasets"


def create_clean_docs(out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1. Plain Text
    with open(out_dir / "clean_sample.txt", "w") as f:
        f.write(
            "This is a totally clean plain text document. It contains vital indexing keywords like 'ALPHA-CLEAN-001'."
        )

    # 2. Markdown
    with open(out_dir / "clean_sample.md", "w") as f:
        f.write(
            "# Clean Markdown\\n\\nHere is a clean markdown file with **bold** text and keywords 'ALPHA-CLEAN-002'."
        )

    # 3. Clean PDF (Vector)
    c = canvas.Canvas(str(out_dir / "clean_vector.pdf"))
    c.drawString(100, 700, "Clean PDF Document.")
    c.drawString(100, 680, "This is natively generated without OCR.")
    c.drawString(100, 660, "Target keyword: ALPHA-CLEAN-003")
    c.save()


def create_scanned_docs(out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)

    # Generate an image with text
    img = Image.new("RGB", (800, 400), color=(255, 255, 255))
    d = ImageDraw.Draw(img)
    d.text((50, 50), "Scanned Invoice 9901", fill=(0, 0, 0))
    d.text((50, 100), "Target keyword: SCANNED-OCR-001", fill=(0, 0, 0))
    img.save(out_dir / "scanned_invoice.png")

    # Generate a PDF wrapping an image (simulating scanned PDF)
    c = canvas.Canvas(str(out_dir / "scanned_wrapper.pdf"))
    c.drawImage(str(out_dir / "scanned_invoice.png"), 50, 400, width=400, height=200)
    c.save()


def create_noisy_images(out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)

    # Generate an image, apply blur and noise
    img = Image.new("RGB", (800, 400), color=(200, 200, 200))
    d = ImageDraw.Draw(img)
    d.text((50, 50), "Poor Quality Fax", fill=(50, 50, 50))
    d.text((50, 150), "Target keyword: NOISY-OCR-001", fill=(30, 30, 30))

    # Blur
    img = img.filter(ImageFilter.GaussianBlur(1.0))
    img.save(out_dir / "noisy_fax.jpg", "JPEG", quality=40)


def create_office_docs(out_dir: Path):
    out_dir.mkdir(parents=True, exist_ok=True)

    # 1. DOCX with Tables
    doc = Document()
    doc.add_heading("Financial Report 2026", 0)
    doc.add_paragraph("Target keyword: OFFICE-WORD-001")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Revenue"
    table.cell(0, 1).text = "$1,000,000"
    table.cell(1, 0).text = "Expenses"
    table.cell(1, 1).text = "$800,000"
    doc.save(str(out_dir / "financial_report.docx"))

    # 2. XLSX with Charts/Tables
    wb = Workbook()
    ws = wb.active
    ws.title = "Q1 Forecast"
    ws.append(["Category", "Amount"])
    ws.append(["Target keyword", "OFFICE-EXCEL-001"])
    ws.append(["Hardware", 50000])
    wb.save(str(out_dir / "q1_forecast.xlsx"))

    # 3. PPTX with charts
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Q1 Presentation"
    subtitle.text = "Target keyword: OFFICE-PPTX-001"
    prs.save(str(out_dir / "q1_presentation.pptx"))


if __name__ == "__main__":
    if DATASET_DIR.exists():
        shutil.rmtree(DATASET_DIR)
    DATASET_DIR.mkdir(parents=True)

    print("Generating Clean Docs...")
    create_clean_docs(DATASET_DIR / "clean")

    print("Generating Scanned Docs...")
    create_scanned_docs(DATASET_DIR / "scanned")

    print("Generating Noisy Images...")
    create_noisy_images(DATASET_DIR / "noisy")

    print("Generating Office Docs...")
    create_office_docs(DATASET_DIR / "office")

    print(f"Dataset generated successfully at {DATASET_DIR}")
