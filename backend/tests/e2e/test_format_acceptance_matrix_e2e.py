from __future__ import annotations

import io
import time
import uuid
from collections.abc import Callable
from typing import Any

# Import purely for creating valid fixture files without saving to disk
import docx  # type: ignore[import-untyped]
import openpyxl  # type: ignore[import-untyped]
import pptx  # type: ignore[import-untyped]
import pytest
from fastapi.testclient import TestClient
from PIL import Image
from pypdf import PdfWriter
from pytest import MonkeyPatch

from app.core.config import Settings
from app.system.services.storage_service import StorageService, StoredObject
from tests.conftest import SeededUser


def _login(client: TestClient, seeded: SeededUser) -> str:
    response = client.post(
        "/api/v1/auth/login",
        headers={"X-Tenant-Id": str(seeded.tenant_id)},
        json={"email": seeded.email, "password": seeded.password},
    )
    assert response.status_code == 200
    return str(response.json()["access_token"])


def _patch_storage(monkeypatch: MonkeyPatch) -> None:
    in_memory_objects: dict[tuple[str, str], bytes] = {}

    def fake_put(
        self: StorageService,
        *,
        tenant_id: uuid.UUID,
        document_id: uuid.UUID,
        filename: str,
        content_type: str,
        payload: bytes,
    ) -> StoredObject:
        object_key = f"{tenant_id}/{document_id}/{filename}"
        in_memory_objects[(self.settings.minio_bucket, object_key)] = payload
        return StoredObject(
            bucket=self.settings.minio_bucket,
            object_key=object_key,
            etag="etag",
            size_bytes=len(payload),
            content_type=content_type,
        )

    def fake_get(self: StorageService, *, bucket: str, object_key: str) -> bytes:
        return in_memory_objects[(bucket, object_key)]

    monkeypatch.setattr(StorageService, "put_bytes", fake_put)
    monkeypatch.setattr(StorageService, "get_bytes", fake_get)


# Fixture generators for various common formats


def _make_pdf(text_content: str) -> bytes:
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    # The pypdf fallback to OCR might require actual pixels or text, but for the E2E matrix,
    # we just need a valid PDF file. In native extract, text is optional if OCR kicks in.
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()


def _make_docx(text_content: str) -> bytes:
    doc = docx.Document()
    doc.add_paragraph(text_content)
    output = io.BytesIO()
    doc.save(output)
    return output.getvalue()


def _make_pptx(text_content: str) -> bytes:
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = text_content
    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()


def _make_xlsx(text_content: str) -> bytes:
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = text_content
    output = io.BytesIO()
    wb.save(output)
    return output.getvalue()


def _make_png() -> bytes:
    image = Image.new("RGB", (100, 100), color=(73, 109, 137))
    output = io.BytesIO()
    image.save(output, format="PNG")
    return output.getvalue()


def wait_for_status(
    client: TestClient, token: str, tenant_id: uuid.UUID, doc_id: str, timeout: int = 10
) -> dict[str, Any]:
    start = time.time()
    while time.time() - start < timeout:
        res = client.get(
            f"/api/v1/documents/{doc_id}/status",
            headers={"Authorization": f"Bearer {token}", "X-Tenant-Id": str(tenant_id)},
        )
        assert res.status_code == 200
        from typing import cast

        data = cast(dict[str, Any], res.json())
        if data["status"] in ("indexed", "failed", "dead_lettered"):
            return data
        time.sleep(0.5)
    raise TimeoutError("Document did not reach terminal status")


@pytest.mark.parametrize(
    "filename, content_type, generator, text_content",
    [
        (
            "test.txt",
            "text/plain",
            lambda t: t.encode("utf-8"),
            "Format test unique phrase TXT.",
        ),
        (
            "test.md",
            "text/markdown",
            lambda t: t.encode("utf-8"),
            "Format test unique phrase MD.",
        ),
        (
            "test.py",
            "text/plain",
            lambda t: t.encode("utf-8"),
            "Format test unique phrase PY.",
        ),
        (
            "test.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            _make_docx,
            "Format test unique phrase DOCX.",
        ),
        (
            "test.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            _make_pptx,
            "Format test unique phrase PPTX.",
        ),
        (
            "test.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            _make_xlsx,
            "Format test unique phrase XLSX.",
        ),
    ],
)
def test_end_to_end_native_formats(
    client: TestClient,
    settings: Settings,
    seed_user: Callable[[str, str, str, tuple[str, ...]], SeededUser],
    monkeypatch: MonkeyPatch,
    filename: str,
    content_type: str,
    generator: Callable[[str], bytes],
    text_content: str,
) -> None:
    _patch_storage(monkeypatch)

    # Ensure test formats are allowed in test config
    if content_type not in settings.upload_allowed_mime_types:
        settings.upload_allowed_mime_types.append(content_type)

    ext = "." + filename.split(".")[-1]
    if ext not in settings.upload_allowed_extensions:
        settings.upload_allowed_extensions.append(ext)

    from app.core.config import get_settings

    client.app.dependency_overrides[get_settings] = lambda: settings  # type: ignore

    seeded = seed_user(
        f"tenant-{filename}",
        f"user-{filename}@example.com",
        "Password!123",
        ("editor",),
    )
    token = _login(client, seeded)

    payload = generator(text_content)

    # 1. Upload
    response = client.post(
        "/api/v1/documents/upload",
        headers={
            "Authorization": f"Bearer {token}",
            "X-Tenant-Id": str(seeded.tenant_id),
            "Idempotency-Key": f"idem-{filename}",
        },
        files={"file": (filename, payload, content_type)},
    )
    if response.status_code != 200:
        pytest.fail(
            f"Upload failed for {filename} with config {settings.upload_allowed_mime_types}: {response.json()}"
        )
    assert response.status_code == 200
    doc_id = response.json()["document_id"]

    # 2. Status reached indexed
    status = wait_for_status(client, token, seeded.tenant_id, doc_id)
    assert (
        status["status"] == "indexed"
    ), f"Failed for {filename}: {status.get('last_error_message')}"
    assert "extraction_method" in status
    assert status["extraction_coverage_score"] is not None

    # 3. Query returns citations
    query = client.post(
        "/api/v1/queries",
        headers={
            "Authorization": f"Bearer {token}",
            "X-Tenant-Id": str(seeded.tenant_id),
        },
        json={"query": "Format test unique phrase", "top_k": 3, "filters": {}},
    )
    assert query.status_code == 200
    q_data = query.json()
    assert len(q_data["citations"]) > 0
    assert any(c["document_id"] == doc_id for c in q_data["citations"])


def test_tenant_isolation_e2e(
    client: TestClient,
    seed_user: Callable[[str, str, str, tuple[str, ...]], SeededUser],
    monkeypatch: MonkeyPatch,
) -> None:
    _patch_storage(monkeypatch)
    tenant_a = seed_user("tenant-a", "user-a@tenant.com", "Password!123", ("editor",))
    tenant_b = seed_user("tenant-b", "user-b@tenant.com", "Password!123", ("editor",))
    token_a = _login(client, tenant_a)
    token_b = _login(client, tenant_b)

    # Tenant A uploads a document
    res = client.post(
        "/api/v1/documents/upload",
        headers={
            "Authorization": f"Bearer {token_a}",
            "X-Tenant-Id": str(tenant_a.tenant_id),
            "Idempotency-Key": "idem-a",
        },
        files={"file": ("secret.txt", b"highly classified tenant a secrets", "text/plain")},
    )
    assert res.status_code == 200
    doc_id_a = res.json()["document_id"]
    status_a = wait_for_status(client, token_a, tenant_a.tenant_id, doc_id_a)
    assert status_a["status"] == "indexed"

    # Tenant B queries the exact keywords
    query_b = client.post(
        "/api/v1/queries",
        headers={
            "Authorization": f"Bearer {token_b}",
            "X-Tenant-Id": str(tenant_b.tenant_id),
        },
        json={"query": "highly classified tenant a secrets", "top_k": 5, "filters": {}},
    )
    assert query_b.status_code == 200
    assert len(query_b.json()["citations"]) == 0  # Tenant B cannot see it
