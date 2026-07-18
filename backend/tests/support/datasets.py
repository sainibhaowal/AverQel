from __future__ import annotations

import shutil
from pathlib import Path

from docx import Document
from openpyxl import Workbook  # type: ignore[import-untyped]
from PIL import Image, ImageDraw
from pptx import Presentation

DATASET_DIR = Path(__file__).parent.parent.parent.parent / "Docs" / "Datasets"


def _create_clean_docs(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    (out_dir / "clean_sample.txt").write_text(
        "This is a totally clean plain text document with keyword ALPHA-CLEAN-001.",
        encoding="utf-8",
    )
    (out_dir / "clean_sample.md").write_text(
        "# Clean Markdown\n\nThis markdown contains keyword ALPHA-CLEAN-002.",
        encoding="utf-8",
    )


def _create_scanned_docs(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (800, 400), color=(255, 255, 255))
    draw = ImageDraw.Draw(image)
    draw.text((50, 50), "Scanned Invoice 9901", fill=(0, 0, 0))
    draw.text((50, 100), "Target keyword: SCANNED-OCR-001", fill=(0, 0, 0))
    image.save(out_dir / "scanned_invoice.png")


def _create_noisy_images(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)
    image = Image.new("RGB", (800, 400), color=(200, 200, 200))
    draw = ImageDraw.Draw(image)
    draw.text((50, 50), "Poor Quality Fax", fill=(50, 50, 50))
    draw.text((50, 150), "Target keyword: NOISY-OCR-001", fill=(30, 30, 30))
    image.save(out_dir / "noisy_fax.jpg", "JPEG", quality=40)


def _create_office_docs(out_dir: Path) -> None:
    out_dir.mkdir(parents=True, exist_ok=True)

    doc = Document()
    doc.add_heading("Financial Report 2026", 0)
    doc.add_paragraph("Target keyword: OFFICE-WORD-001")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Revenue"
    table.cell(0, 1).text = "$1,000,000"
    table.cell(1, 0).text = "Expenses"
    table.cell(1, 1).text = "$800,000"
    doc.save(str(out_dir / "financial_report.docx"))

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Q1 Forecast"
    worksheet.append(["Category", "Amount"])
    worksheet.append(["Target keyword", "OFFICE-EXCEL-001"])
    worksheet.append(["Hardware", 50000])
    workbook.save(str(out_dir / "q1_forecast.xlsx"))

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Q1 Presentation"
    slide.placeholders[1].text = "Target keyword: OFFICE-PPTX-001"
    presentation.save(str(out_dir / "q1_presentation.pptx"))


def ensure_test_datasets() -> Path:
    if DATASET_DIR.exists():
        return DATASET_DIR

    DATASET_DIR.mkdir(parents=True, exist_ok=True)
    _create_clean_docs(DATASET_DIR / "clean")
    _create_scanned_docs(DATASET_DIR / "scanned")
    _create_noisy_images(DATASET_DIR / "noisy")
    _create_office_docs(DATASET_DIR / "office")
    return DATASET_DIR


def reset_test_datasets() -> Path:
    if DATASET_DIR.exists():
        shutil.rmtree(DATASET_DIR)
    return ensure_test_datasets()
